"""
processors/pptx_renderer.py
────────────────────────────────────────────────────────────────────
PPTX → PNG renderer (python-pptx + Pillow)

PowerPoint COM veya LibreOffice yoksa fallback olarak çalışır.
Her slaydı doğrudan shape ağacından Pillow Image üzerine çizer:
  - PICTURE shape  → gömülü blob'u doğru EMU konumuna paste
  - TEXT_BOX/AUTO  → metin + arka plan kutusu
  - LINE           → çizgi
  - AUTO_SHAPE     → dikdörtgen çerçeve
────────────────────────────────────────────────────────────────────
"""

import io
import os
from typing import Optional

from PIL import Image, ImageDraw, ImageFont

from core.logger import get_logger

logger = get_logger("processors.pptx_renderer")

_DPI = 144  # Render çözünürlüğü (1 inch = 144 px)
_EMU_PER_INCH = 914400


def _emu_to_px(emu: int, dpi: int = _DPI) -> int:
    return int(round(emu * dpi / _EMU_PER_INCH))


def _load_font(size_px: int) -> ImageFont.FreeTypeFont:
    candidates = [
        "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/segoeui.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    for path in candidates:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size_px)
            except Exception:
                continue
    return ImageFont.load_default()


def _draw_picture(canvas: Image.Image, shape, scale_x: float, scale_y: float) -> None:
    try:
        blob = shape.image.blob
        img = Image.open(io.BytesIO(blob)).convert("RGBA")

        x = _emu_to_px(shape.left or 0)
        y = _emu_to_px(shape.top or 0)
        w = _emu_to_px(shape.width or 0)
        h = _emu_to_px(shape.height or 0)

        if w <= 0 or h <= 0:
            return

        img = img.resize((w, h), Image.LANCZOS)
        canvas.paste(img, (x, y), img)
    except Exception as e:
        logger.debug(f"[Renderer] Picture çizilemedi: {e}")


def _draw_text(canvas: Image.Image, draw: ImageDraw.ImageDraw, shape) -> None:
    try:
        if not shape.has_text_frame:
            return
        text = shape.text_frame.text or ""
        if not text.strip():
            return

        x = _emu_to_px(shape.left or 0)
        y = _emu_to_px(shape.top or 0)
        w = _emu_to_px(shape.width or 0)
        h = _emu_to_px(shape.height or 0)

        if w <= 0 or h <= 0:
            return

        # Hafif arka plan (kenarlığı belli olsun)
        draw.rectangle([x, y, x + w, y + h], fill=(255, 255, 255, 255), outline=(200, 200, 200, 255))

        # Font boyutu kutuya göre
        font_size = max(10, min(18, h // 6))
        font = _load_font(font_size)

        # Metni satırlara böl ve sığdır
        pad = 4
        max_chars_per_line = max(8, (w - 2 * pad) // (font_size // 2))
        lines: list[str] = []
        for raw_line in text.split("\n"):
            words = raw_line.split()
            if not words:
                lines.append("")
                continue
            cur = ""
            for word in words:
                trial = (cur + " " + word).strip()
                if len(trial) > max_chars_per_line:
                    if cur:
                        lines.append(cur)
                    cur = word
                else:
                    cur = trial
            if cur:
                lines.append(cur)

        line_height = font_size + 2
        max_lines = max(1, (h - 2 * pad) // line_height)
        for i, line in enumerate(lines[:max_lines]):
            draw.text((x + pad, y + pad + i * line_height), line, fill=(40, 40, 40), font=font)
    except Exception as e:
        logger.debug(f"[Renderer] Text çizilemedi: {e}")


def _draw_line(draw: ImageDraw.ImageDraw, shape) -> None:
    try:
        x1 = _emu_to_px(shape.left or 0)
        y1 = _emu_to_px(shape.top or 0)
        x2 = x1 + _emu_to_px(shape.width or 0)
        y2 = y1 + _emu_to_px(shape.height or 0)
        draw.line([(x1, y1), (x2, y2)], fill=(80, 80, 80), width=2)
    except Exception as e:
        logger.debug(f"[Renderer] Line çizilemedi: {e}")


def _draw_auto_shape(draw: ImageDraw.ImageDraw, shape) -> None:
    try:
        x = _emu_to_px(shape.left or 0)
        y = _emu_to_px(shape.top or 0)
        w = _emu_to_px(shape.width or 0)
        h = _emu_to_px(shape.height or 0)
        if w <= 0 or h <= 0:
            return
        draw.rectangle([x, y, x + w, y + h], outline=(120, 120, 120), width=1)
    except Exception as e:
        logger.debug(f"[Renderer] AutoShape çizilemedi: {e}")


def _render_slide_to_image(slide, slide_w_emu: int, slide_h_emu: int) -> Optional[Image.Image]:
    """Tek slaydı Pillow Image olarak döner (diske yazmadan)."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        canvas_w = _emu_to_px(slide_w_emu)
        canvas_h = _emu_to_px(slide_h_emu)
        canvas = Image.new("RGBA", (canvas_w, canvas_h), (255, 255, 255, 255))
        draw = ImageDraw.Draw(canvas)

        for shape in slide.shapes:
            try:
                stype = shape.shape_type
                if stype == MSO_SHAPE_TYPE.PICTURE:
                    _draw_picture(canvas, shape, 1.0, 1.0)
                elif stype == MSO_SHAPE_TYPE.LINE:
                    _draw_line(draw, shape)
                elif shape.has_text_frame and (shape.text_frame.text or "").strip():
                    _draw_text(canvas, draw, shape)
                elif stype == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    _draw_auto_shape(draw, shape)
            except Exception as e:
                logger.debug(f"[Renderer] Shape atlandı: {e}")
                continue

        return canvas.convert("RGB")
    except Exception as e:
        logger.warning(f"[Renderer] Slayt render başarısız: {e}")
        return None


def render_pptx_to_pngs(pptx_path: str, output_dir: str) -> list[str]:
    """PPTX'in tüm slaytlarını PNG olarak üretir, üretilen dosya yollarını döner."""
    from pptx import Presentation

    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    sw = prs.slide_width
    sh = prs.slide_height

    produced: list[str] = []
    for idx, slide in enumerate(prs.slides):
        out_path = os.path.join(output_dir, f"page_{idx + 1}.png")
        img = _render_slide_to_image(slide, sw, sh)
        if img is not None:
            img.save(out_path, "PNG", optimize=True)
            produced.append(out_path)
        else:
            logger.warning(f"[Renderer] Slayt {idx + 1} üretilemedi: {pptx_path}")

    return produced


def render_pptx_to_pdf(pptx_path: str, pdf_path: str) -> Optional[str]:
    """
    PPTX'in tüm slaytlarını render edip tek bir PDF dosyasına birleştirir.
    Pillow ile raster PDF üretir (vektör değil — slayt başına bir görsel).
    """
    from pptx import Presentation

    os.makedirs(os.path.dirname(os.path.abspath(pdf_path)), exist_ok=True)
    prs = Presentation(pptx_path)
    sw = prs.slide_width
    sh = prs.slide_height

    images: list[Image.Image] = []
    for idx, slide in enumerate(prs.slides):
        img = _render_slide_to_image(slide, sw, sh)
        if img is not None:
            images.append(img)
        else:
            logger.warning(f"[Renderer] PDF için slayt {idx + 1} üretilemedi")

    if not images:
        logger.warning(f"[Renderer] Hiçbir slayt üretilemedi, PDF oluşturulmadı: {pptx_path}")
        return None

    try:
        first, rest = images[0], images[1:]
        first.save(
            pdf_path,
            "PDF",
            resolution=_DPI,
            save_all=True,
            append_images=rest,
        )
        return pdf_path
    except Exception as e:
        logger.warning(f"[Renderer] PDF kaydetme hatası: {e}")
        return None


def render_pptx_full(pptx_path: str, output_dir: str) -> dict:
    """
    PPTX'i tek geçişte render eder: hem PNG'leri hem birleşik PDF üretir.
    Aynı Pillow Image'ları kullandığından çift render maliyeti yok.
    Döner: {"pdf": <path|None>, "pngs": [<path>, ...]}
    """
    from pptx import Presentation

    os.makedirs(output_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    sw = prs.slide_width
    sh = prs.slide_height

    images: list[Image.Image] = []
    pngs: list[str] = []

    for idx, slide in enumerate(prs.slides):
        img = _render_slide_to_image(slide, sw, sh)
        if img is None:
            logger.warning(f"[Renderer] Slayt {idx + 1} üretilemedi")
            continue
        png_path = os.path.join(output_dir, f"page_{idx + 1}.png")
        img.save(png_path, "PNG", optimize=True)
        pngs.append(png_path)
        images.append(img)

    pdf_path: Optional[str] = None
    if images:
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        pdf_path = os.path.join(os.path.dirname(pptx_path), f"{base_name}.pdf")
        try:
            first, rest = images[0], images[1:]
            first.save(
                pdf_path,
                "PDF",
                resolution=_DPI,
                save_all=True,
                append_images=rest,
            )
        except Exception as e:
            logger.warning(f"[Renderer] PDF üretilemedi: {e}")
            pdf_path = None

    return {"pdf": pdf_path, "pngs": pngs}
