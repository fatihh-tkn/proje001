"""
processors/pptx_processor.py  ─ v2
────────────────────────────────────────────────────────────────────
PPTX → RAG chunk üretici (Per-Shape Hibrit Mimari)

  AŞAMA 1 ─ Per-Shape Chunking
    Her text shape → kendi bbox'ıyla ayrı chunk
    type: pptx_title / pptx_callout / pptx_body / pptx_notes

  AŞAMA 2 ─ Connector/Arrow Linking
    Ok bağlantıları tespit edilir, callout → target_bbox ilişkisi
    kurulur; popup görseli callout değil hedef UI alanını gösterir

  AŞAMA 3 ─ Vision OCR (use_vision=True)
    Ekran görüntüsü shape'leri Gemini'ye gönderilir, form alanı
    isimleri ve değerleri text'e dönüştürülür
    type: pptx_screenshot_vision
────────────────────────────────────────────────────────────────────
"""

import os
import uuid
import shutil
import subprocess

try:
    import fitz
    _FITZ_AVAILABLE = True
except ImportError:
    _FITZ_AVAILABLE = False

_CHUNK_SIZE = 3000
_MIN_TEXT   = 10


# ── PPT → PDF Dönüştürücü ────────────────────────────────────────────

def convert_pptx_to_pdf(pptx_path: str, output_dir: str) -> str | None:
    base_name  = os.path.splitext(os.path.basename(pptx_path))[0]
    target_pdf = os.path.join(output_dir, f"{base_name}.pdf")

    if os.path.exists(target_pdf) and os.path.getsize(target_pdf) > 0:
        return target_pdf

    os.makedirs(output_dir, exist_ok=True)

    try:
        import comtypes.client
        pptx_abs = os.path.abspath(pptx_path)
        pdf_abs  = os.path.abspath(target_pdf)
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        deck = powerpoint.Presentations.Open(pptx_abs, ReadOnly=True, Untitled=True, WithWindow=False)
        deck.SaveAs(pdf_abs, 32)
        deck.Close()
        powerpoint.Quit()
        print(f"[PPTX->PDF] pywin32 (COM) başarılı: {pdf_abs}")
        return pdf_abs
    except Exception as e:
        print(f"[PPTX->PDF] pywin32 kullanılamıyor ({e}), LibreOffice deneniyor...")

    soffice = shutil.which("soffice") or shutil.which("soffice.bin")
    if soffice:
        try:
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", output_dir, pptx_path],
                capture_output=True, text=True, timeout=120,
            )
            expected = os.path.join(output_dir, f"{base_name}.pdf")
            if result.returncode == 0 and os.path.exists(expected):
                print(f"[PPTX->PDF] LibreOffice başarılı: {expected}")
                return expected
            else:
                print(f"[PPTX->PDF] LibreOffice başarısız: {result.stderr[:200]}")
        except Exception as e:
            print(f"[PPTX->PDF] LibreOffice hatası: {e}")
    else:
        print("[PPTX->PDF] Dönüştürücü bulunamadı — yalnızca metin chunking.")

    return None


# ── Ana Parser ───────────────────────────────────────────────────────

def parse_pptx(
    file_path: str,
    original_name: str | None = None,
    use_vision: bool = False,
) -> list[dict]:
    file_basename = original_name or os.path.basename(file_path)
    output_dir    = os.path.dirname(file_path)
    pdf_path      = convert_pptx_to_pdf(file_path, output_dir)
    print(f"[PPTX] Per-shape parse (PDF: {bool(pdf_path)}, Vision: {use_vision}): {file_basename}")
    return _hybrid_parse_pptx(file_path, file_basename, pdf_path, use_vision=use_vision)


# ── Çekirdek İşleyici ────────────────────────────────────────────────

def _hybrid_parse_pptx(
    file_path: str,
    file_basename: str,
    pdf_path: str | None,
    use_vision: bool = False,
) -> list[dict]:
    chunks: list[dict] = []

    try:
        from pptx import Presentation
    except ImportError:
        return [_error_chunk(file_basename, "python-pptx kurulu değil. pip install python-pptx")]

    try:
        prs          = Presentation(file_path)
        total_slides = len(prs.slides)
        slide_width  = prs.slide_width
        slide_height = prs.slide_height
    except Exception as e:
        return [_error_chunk(file_basename, f"PPTX okunamadı: {e}")]

    # ── Fitz / PNG hazırlığı ──────────────────────────────────────────
    doc = None
    image_dir = ""
    if _FITZ_AVAILABLE and pdf_path and os.path.exists(pdf_path):
        try:
            doc = fitz.open(pdf_path)
            base_name = os.path.splitext(file_basename)[0]
            image_dir = os.path.join(os.path.dirname(file_path), f"images_{base_name}")
            os.makedirs(image_dir, exist_ok=True)
        except Exception as e:
            print(f"[PPTX] Fitz başlatılamadı: {e}")
            doc = None

    slide_titles: list[str] = []

    # ── Her slaytı işle ───────────────────────────────────────────────
    for slide_idx, slide in enumerate(prs.slides):
        title       = _get_title(slide)
        table_texts = _get_table_texts(slide)
        notes_text  = _get_notes(slide)
        slide_titles.append(title or f"Slayt {slide_idx + 1}")

        header = f"[{file_basename} | Slayt {slide_idx+1}/{total_slides}]"

        # ── Slayt PNG çıkart ──────────────────────────────────────────
        slide_image_path = ""
        if _FITZ_AVAILABLE and doc and slide_idx < len(doc):
            try:
                page = doc.load_page(slide_idx)
                pix  = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                fname = f"page_{slide_idx + 1}.png"
                img_path_abs = os.path.abspath(os.path.join(image_dir, fname))
                pix.save(img_path_abs)
                slide_image_path = img_path_abs
            except Exception as e:
                print(f"[PPTX] PNG kaydedilemedi slayt {slide_idx+1}: {e}")

        # ── AŞAMA 1: Per-shape shapes ─────────────────────────────────
        shapes_data   = _extract_shapes_individual(slide)
        # ── AŞAMA 2: Connector mapping ────────────────────────────────
        connector_map = _detect_connectors(slide)   # {shape_id: target_bbox_str}

        base_meta = {
            "page":              slide_idx + 1,
            "source":            file_basename,
            "total_pages":       total_slides,
            "file_type":         "pptx",
            "image_path":        slide_image_path,
            "slide_emu_w":       int(slide_width),
            "slide_emu_h":       int(slide_height),
            "slide_chunk_group": f"{file_basename}::slide{slide_idx + 1}",
        }

        # ── Per-shape chunk'lar ───────────────────────────────────────
        shape_count = 0
        for sd in shapes_data:
            text = sd["text"]
            if len(text.strip()) < _MIN_TEXT and sd["shape_type"] != "title":
                continue

            target_bbox = connector_map.get(sd["shape_id"], "")

            chunk_text = header + "\n"
            if title and sd["shape_type"] != "title":
                chunk_text += f"BAŞLIK: {title}\n"
            chunk_text += text

            shape_count += 1
            chunks.append({
                "id":   str(uuid.uuid4()),
                "text": chunk_text,
                "metadata": {
                    **base_meta,
                    "chunk_index": shape_count,
                    "type":        f"pptx_{sd['shape_type']}",
                    "bbox":        sd["bbox"],
                    "target_bbox": target_bbox,
                    "shape_id":    str(sd["shape_id"]),
                }
            })

        # ── Tablo / Notlar chunk'ı ────────────────────────────────────
        extra_parts = []
        if table_texts:
            extra_parts.append("TABLO:\n" + "\n".join(table_texts))
        if notes_text:
            extra_parts.append(f"KONUŞMACI NOTLARI:\n{notes_text}")

        if extra_parts:
            shape_count += 1
            chunks.append({
                "id":   str(uuid.uuid4()),
                "text": f"{header}\nBAŞLIK: {title or ''}\n\n" + "\n\n".join(extra_parts),
                "metadata": {
                    **base_meta,
                    "chunk_index": shape_count,
                    "type":        "pptx_notes",
                    "bbox":        "0,0,0,0",
                    "target_bbox": "",
                    "shape_id":    "",
                }
            })

        # ── AŞAMA 3: Vision OCR — ekran görüntüsü shape'leri ─────────
        if use_vision and slide_image_path:
            picture_shapes = _extract_picture_shapes(slide)
            for pic_idx, pic in enumerate(picture_shapes):
                cropped = _crop_shape_image(
                    slide_image_path, pic["bbox"],
                    int(slide_width), int(slide_height),
                    image_dir, slide_idx, pic["shape_id"],
                )
                if not cropped:
                    continue
                vision_text = _ask_gemini_for_ui(cropped, context=title or "")
                if not vision_text or vision_text.startswith("["):
                    continue
                shape_count += 1
                chunks.append({
                    "id":   str(uuid.uuid4()),
                    "text": f"{header}\nBAŞLIK: {title or ''}\n\nEKRAN GÖRÜNTÜSİ ANALİZİ:\n{vision_text}",
                    "metadata": {
                        **base_meta,
                        "chunk_index": shape_count,
                        "type":        "pptx_screenshot_vision",
                        "bbox":        pic["bbox"],
                        "target_bbox": "",
                        "shape_id":    str(pic["shape_id"]),
                    }
                })

        # Hiç shape yoksa başlıkla minimal chunk ekle
        if shape_count == 0:
            if title or table_texts:
                fallback_text = header
                if title:
                    fallback_text += f"\nBAŞLIK: {title}"
                if table_texts:
                    fallback_text += "\nTABLO:\n" + "\n".join(table_texts)
                chunks.append({
                    "id":   str(uuid.uuid4()),
                    "text": fallback_text,
                    "metadata": {
                        **base_meta,
                        "chunk_index": 1,
                        "type":        "pptx_slide",
                        "bbox":        "0,0,0,0",
                        "target_bbox": "",
                        "shape_id":    "",
                    }
                })

    if doc:
        doc.close()

    # ── Özet (İçindekiler) chunk'ı ────────────────────────────────────
    non_empty = [f"  {i+1}. {t}" for i, t in enumerate(slide_titles) if t.strip()]
    if non_empty:
        chunks.insert(0, {
            "id":   str(uuid.uuid4()),
            "text": (
                f"SUNUM: {file_basename}\n"
                f"TOPLAM SLAYT: {total_slides}\n\n"
                "── SLAYT BAŞLIKLARI (İÇİNDEKİLER) ──\n"
                + "\n".join(non_empty)
            ),
            "metadata": {
                "page":              0,
                "chunk_index":       0,
                "source":            file_basename,
                "type":              "pptx_summary",
                "total_pages":       total_slides,
                "file_type":         "pptx",
                "image_path":        "",
                "bbox":              "0,0,0,0",
                "target_bbox":       "",
                "shape_id":          "",
                "slide_emu_w":       int(slide_width),
                "slide_emu_h":       int(slide_height),
                "slide_chunk_group": f"{file_basename}::summary",
            }
        })

    if not chunks:
        chunks.append(_error_chunk(file_basename, "Slayt içeriği bulunamadı."))

    return chunks


# ── AŞAMA 1: Per-shape çıkarım ───────────────────────────────────────

def _extract_shapes_individual(slide) -> list[dict]:
    """Her text shape'i kendi bbox'ıyla döner, Y+X'e göre sıralı."""
    title_shape_id = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            title_shape_id = ph.shape_id
            break

    result = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue

        left = shape.left  or 0
        top  = shape.top   or 0
        w    = shape.width or 0
        h    = shape.height or 0
        bbox = f"{int(left)},{int(top)},{int(left + w)},{int(top + h)}"

        if shape.shape_id == title_shape_id:
            shape_type = "title"
        elif shape.is_placeholder:
            shape_type = "body"
        else:
            shape_type = _classify_shape(shape)

        result.append({
            "shape_id":   shape.shape_id,
            "text":       text,
            "bbox":       bbox,
            "shape_type": shape_type,
            "_top":       int(top),
            "_left":      int(left),
        })

    result.sort(key=lambda x: (x["_top"], x["_left"]))
    return result


def _classify_shape(shape) -> str:
    """Şekli callout / body olarak sınıflandırır."""
    name = (shape.name or "").lower()
    if any(k in name for k in ("callout", "açıklama", "note", "balloon", "caption")):
        return "callout"

    # python-pptx auto_shape_type: callout aralığı 43-66
    try:
        ast = shape.auto_shape_type
        if ast is not None and 43 <= int(ast) <= 66:
            return "callout"
    except Exception:
        pass

    return "callout"  # varsayılan: placeholder dışı text box = callout


# ── AŞAMA 2: Connector tespiti ───────────────────────────────────────

def _detect_connectors(slide) -> dict[int, str]:
    """
    Slayttaki ok/bağlantı şekillerini tarar.
    Döner: {callout_shape_id: target_bbox_str}
    """
    from pptx.oxml.ns import qn

    # shape_id → shape hızlı erişim
    shape_map: dict[int, object] = {s.shape_id: s for s in slide.shapes}

    result: dict[int, str] = {}

    for shape in slide.shapes:
        # Sadece connector (cxnSp) elementleri
        if shape._element.tag != qn("p:cxnSp"):
            continue

        try:
            cf = shape.connector_format
            begin_shape = cf.begin_connected_shape
            end_shape   = cf.end_connected_shape
        except Exception:
            continue

        if begin_shape is None or end_shape is None:
            continue

        begin_text = begin_shape.has_text_frame and begin_shape.text_frame.text.strip()
        end_text   = end_shape.has_text_frame   and end_shape.text_frame.text.strip()

        # Metinli taraf = callout, diğer taraf = hedef
        if begin_text and not end_text:
            callout_sh, target_sh = begin_shape, end_shape
        elif end_text and not begin_text:
            callout_sh, target_sh = end_shape, begin_shape
        else:
            continue

        tl = target_sh.left  or 0
        tt = target_sh.top   or 0
        tw = target_sh.width or 0
        th = target_sh.height or 0
        target_bbox = f"{int(tl)},{int(tt)},{int(tl+tw)},{int(tt+th)}"
        result[callout_sh.shape_id] = target_bbox

    return result


# ── AŞAMA 3: Görsel shape tespiti ve Vision OCR ──────────────────────

def _extract_picture_shapes(slide) -> list[dict]:
    """Slayttaki resim/ekran görüntüsü shape'lerini döner."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        pic_types = {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE}
    except Exception:
        pic_types = set()

    result = []
    for shape in slide.shapes:
        is_pic = False
        try:
            if shape.shape_type in pic_types:
                is_pic = True
        except Exception:
            pass

        if not is_pic:
            # XML tag ile kontrol (p:pic)
            try:
                from pptx.oxml.ns import qn
                if shape._element.tag == qn("p:pic"):
                    is_pic = True
            except Exception:
                pass

        if is_pic:
            left = shape.left  or 0
            top  = shape.top   or 0
            w    = shape.width or 0
            h    = shape.height or 0
            result.append({
                "shape_id": shape.shape_id,
                "bbox":     f"{int(left)},{int(top)},{int(left+w)},{int(top+h)}",
            })

    return result


def _crop_shape_image(
    slide_image_path: str,
    bbox_emu: str,
    slide_emu_w: int,
    slide_emu_h: int,
    image_dir: str,
    slide_idx: int,
    shape_id: int,
) -> str | None:
    """Slide PNG'den belirtilen shape alanını kırpar, path döner."""
    try:
        from PIL import Image
        l, t, r, b = [float(x) for x in bbox_emu.split(",")]
        with Image.open(slide_image_path) as img:
            iw, ih = img.size
            px_l = max(0,  int((l / slide_emu_w) * iw) - 8)
            px_t = max(0,  int((t / slide_emu_h) * ih) - 8)
            px_r = min(iw, int((r / slide_emu_w) * iw) + 8)
            px_b = min(ih, int((b / slide_emu_h) * ih) + 8)
            if px_r <= px_l or px_b <= px_t:
                return None
            cropped = img.crop((px_l, px_t, px_r, px_b))
            out = os.path.abspath(
                os.path.join(image_dir, f"page_{slide_idx+1}_s{shape_id}.png")
            )
            cropped.save(out)
            return out
    except Exception as e:
        print(f"[PPTX] Shape crop hatası: {e}")
        return None


def _ask_gemini_for_ui(image_path: str, context: str = "") -> str:
    """Ekran görüntüsündeki form alanlarını Gemini ile çıkarır."""
    try:
        from core.settings import settings
        if not settings.GEMINI_API_KEY:
            return ""
        import google.generativeai as genai
        from PIL import Image
        genai.configure(api_key=settings.GEMINI_API_KEY)
        model = genai.GenerativeModel("gemini-1.5-pro")

        img = Image.open(image_path)
        ctx = f"Slayt başlığı: '{context}'. " if context else ""
        prompt = (
            f"{ctx}Bu SAP/ERP uygulama ekranındaki form alanlarını, etiketlerini ve "
            "görünen değerleri listele. Her alan için tek satır: 'Alan Adı: değer'. "
            "Başka yorum veya açıklama ekleme."
        )
        response = model.generate_content([prompt, img])
        return response.text.strip()
    except Exception as e:
        print(f"[PPTX Vision] Hata: {e}")
        return ""


# ── Yardımcı fonksiyonlar ────────────────────────────────────────────

def _get_title(slide) -> str:
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0 and ph.has_text_frame:
            t = ph.text_frame.text.strip()
            if t:
                return t
    max_size, best = 0, ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    size = run.font.size or 0
                    text = run.text.strip()
                    if size > max_size and text:
                        max_size, best = size, text
    return best


def _get_table_texts(slide) -> list[str]:
    rows = []
    for shape in slide.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                cells = [c.text.strip() for c in row.cells]
                rows.append(" | ".join(c for c in cells if c))
    return rows


def _get_notes(slide) -> str:
    try:
        tf = slide.notes_slide.notes_text_frame
        if tf:
            lines = [ln for ln in tf.text.strip().splitlines() if ln.strip()]
            return "\n".join(lines)
    except Exception:
        pass
    return ""


def _split(text: str) -> list[str]:
    chunks, start = [], 0
    while start < len(text):
        end = start + _CHUNK_SIZE
        if end < len(text):
            pos = text.rfind("\n", start, end)
            if pos > start:
                end = pos
        part = text[start:end].strip()
        if part:
            chunks.append(part)
        start = end - 400
        if start >= len(text):
            break
    return chunks


def _error_chunk(source: str, msg: str) -> dict:
    return {
        "id":       f"error-{uuid.uuid4()}",
        "text":     f"[{source}] {msg}",
        "metadata": {"source": source, "error": msg, "type": "pptx_error",
                     "page": 0, "chunk_index": 0, "file_type": "pptx",
                     "image_path": "", "bbox": "0,0,0,0", "target_bbox": ""},
    }
